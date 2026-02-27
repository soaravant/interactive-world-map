#!/usr/bin/env python3
"""
Extract images from PowerPoint slides 4-26 for the interactive world map.
Each slide = one mission country. Extracts:
- Group slideshow images (from AUTO_SHAPE placeholders inside GROUP shapes)
- Main photo(s) from the slide
Saves to images/<country_slug>/ directories.
"""
from pptx import Presentation
from pptx.util import Emu
import os
import re
import unicodedata

# Transliteration map for Greek -> Latin for folder names
GREEK_TO_LATIN = {
    'α': 'a', 'β': 'v', 'γ': 'g', 'δ': 'd', 'ε': 'e', 'ζ': 'z', 'η': 'i',
    'θ': 'th', 'ι': 'i', 'κ': 'k', 'λ': 'l', 'μ': 'm', 'ν': 'n', 'ξ': 'x',
    'ο': 'o', 'π': 'p', 'ρ': 'r', 'σ': 's', 'ς': 's', 'τ': 't', 'υ': 'y',
    'φ': 'f', 'χ': 'ch', 'ψ': 'ps', 'ω': 'o',
    'ά': 'a', 'έ': 'e', 'ή': 'i', 'ί': 'i', 'ό': 'o', 'ύ': 'y', 'ώ': 'o',
    'ϊ': 'i', 'ϋ': 'y', 'ΐ': 'i', 'ΰ': 'y',
}

def greek_to_slug(name):
    """Convert Greek name to a URL-safe slug."""
    result = []
    for ch in name.lower():
        if ch in GREEK_TO_LATIN:
            result.append(GREEK_TO_LATIN[ch])
        elif ch.isascii() and ch.isalnum():
            result.append(ch)
        elif ch in (' ', '-', '_'):
            result.append('_')
    slug = ''.join(result)
    slug = re.sub(r'_+', '_', slug).strip('_')
    return slug

BASE_DIR = "/Users/sotiris/Documents/Github Projects/interactive-world-map"
PPTX_PATH = os.path.join(BASE_DIR, "Χάρτης Εξωτερικής Ιεραποστολής.pptx")
IMAGES_DIR = os.path.join(BASE_DIR, "images")

prs = Presentation(PPTX_PATH)

for slide_idx, slide in enumerate(prs.slides):
    slide_num = slide_idx + 1
    if slide_num < 4:
        continue

    # 1. Find country name
    country_name = None
    for shape in slide.shapes:
        if shape.shape_type == 1 and hasattr(shape, "text") and shape.text.strip():
            country_name = shape.text.strip()
            break

    if not country_name:
        print(f"Slide {slide_num}: No country name found, skipping")
        continue

    slug = greek_to_slug(country_name)
    country_dir = os.path.join(IMAGES_DIR, slug)
    os.makedirs(country_dir, exist_ok=True)

    print(f"\nSlide {slide_num}: {country_name} -> {slug}/")

    photo_idx = 1

    # 2. Extract group images (the 3 slideshow photos in AUTO_SHAPE placeholders)
    for shape in slide.shapes:
        if shape.shape_type == 6:  # GROUP
            for sub in shape.shapes:
                if sub.shape_type == 1:  # AUTO_SHAPE (image placeholder)
                    # Check if this AUTO_SHAPE has a fill with an image
                    try:
                        fill = sub.fill
                        if fill.type is not None:
                            # Try to get the image from the fill
                            from pptx.oxml.ns import qn
                            blipFill = sub._element.find('.//' + qn('a:blip'))
                            if blipFill is not None:
                                rId = blipFill.get(qn('r:embed'))
                                if rId:
                                    rel = slide.part.rels[rId]
                                    image_blob = rel.target_part.blob
                                    content_type = rel.target_part.content_type
                                    ext = content_type.split('/')[-1]
                                    if ext == 'jpeg':
                                        ext = 'jpg'
                                    filename = f"photo_{photo_idx}.{ext}"
                                    filepath = os.path.join(country_dir, filename)
                                    with open(filepath, 'wb') as f:
                                        f.write(image_blob)
                                    print(f"  Saved group image: {filename} ({len(image_blob)} bytes)")
                                    photo_idx += 1
                    except Exception as e:
                        pass

    # 3. Extract main photo(s) from slide-level PICTURE shapes
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            blob = shape.image.blob
            w = shape.width
            h = shape.height
            ct = shape.image.content_type

            # Skip background (full slide width > 10M EMU)
            if w > 10000000 and h > 6000000:
                continue

            # Skip tiny icon (the small logo at bottom-right, ~765666 wide)
            if w < 1000000 and h < 1000000:
                continue

            # Skip stats image (Εικόνα)
            if 'Εικόνα' in shape.name:
                continue

            ext = ct.split('/')[-1]
            if ext == 'jpeg':
                ext = 'jpg'
            filename = f"photo_{photo_idx}.{ext}"
            filepath = os.path.join(country_dir, filename)
            with open(filepath, 'wb') as f:
                f.write(blob)
            print(f"  Saved main photo: {filename} ({len(blob)} bytes)")
            photo_idx += 1

    print(f"  Total photos saved: {photo_idx - 1}")

print("\n\nDone! All images extracted.")
